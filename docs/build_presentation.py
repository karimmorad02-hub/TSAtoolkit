#!/usr/bin/env python3
"""
Build CTO-level PPTX presentation for aic_ts_suite toolkit.
Uses python-pptx with custom drawing primitives for diagrams.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import math

# ── Brand colours ──────────────────────────────────────────────
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy
C_BG       = RGBColor(0x0F, 0x0F, 0x1A)   # Slide bg dark
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xE0, 0xE0, 0xE8)
C_MUTED    = RGBColor(0x9A, 0x9A, 0xAE)
C_ACCENT1  = RGBColor(0x4E, 0xA8, 0xF6)   # Blue
C_ACCENT2  = RGBColor(0xA8, 0x55, 0xF7)   # Purple
C_ACCENT3  = RGBColor(0x22, 0xD3, 0xEE)   # Cyan
C_ACCENT4  = RGBColor(0x34, 0xD3, 0x99)   # Green
C_ACCENT5  = RGBColor(0xFB, 0xBF, 0x24)   # Amber
C_ACCENT6  = RGBColor(0xF4, 0x72, 0xB6)   # Pink
C_RED      = RGBColor(0xEF, 0x44, 0x44)   # Red for risk
C_CARD_BG  = RGBColor(0x1E, 0x1E, 0x36)   # Card background
C_BORDER   = RGBColor(0x33, 0x33, 0x50)   # Card border

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helper Functions ───────────────────────────────────────────

def set_slide_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=C_LIGHT):
    """lines = [(text, size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text = line[0]
        size = line[1] if len(line) > 1 else default_size
        color = line[2] if len(line) > 2 else default_color
        bold = line[3] if len(line) > 3 else False
        align = line[4] if len(line) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, title, body_lines, accent_color=C_ACCENT1):
    """Card with accent top border"""
    # Accent line on top
    add_shape(slide, left, top, width, Pt(4), fill_color=accent_color, shape_type=MSO_SHAPE.RECTANGLE)
    # Card body
    add_shape(slide, left, top + Pt(4), width, height - Pt(4), fill_color=C_CARD_BG, border_color=C_BORDER)
    # Title
    add_text(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
             title, font_size=13, color=accent_color, bold=True)
    # Body
    y = top + Inches(0.45)
    for line in body_lines:
        add_text(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.22),
                 line, font_size=10, color=C_LIGHT)
        y += Inches(0.22)

def add_arrow(slide, start_left, start_top, end_left, end_top, color=C_ACCENT1, width=Pt(2)):
    """Add a connector arrow between two points"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_pill(slide, left, top, width, height, text, fill_color=C_ACCENT1, text_color=C_WHITE, font_size=10):
    shape = add_shape(slide, left, top, width, height, fill_color=fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    return shape

def slide_title_bar(slide, title, subtitle=None):
    """Add consistent title bar to all slides"""
    # Title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.55),
             title, font_size=28, color=C_WHITE, bold=True)
    # Accent line
    add_shape(slide, Inches(0.6), Inches(0.85), Inches(1.2), Pt(3),
              fill_color=C_ACCENT1, shape_type=MSO_SHAPE.RECTANGLE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.35),
                 subtitle, font_size=14, color=C_MUTED)

def add_excalidraw_badge(slide, url, left=None, top=None):
    """Add an Excalidraw link badge"""
    if left is None:
        left = Inches(10.5)
    if top is None:
        top = Inches(0.35)
    shape = add_pill(slide, left, top, Inches(2.2), Inches(0.35),
                     "View on Excalidraw", fill_color=C_ACCENT2, font_size=9)
    shape.click_action.hyperlink.address = url

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative gradient bar at top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=C_ACCENT1, shape_type=MSO_SHAPE.RECTANGLE)

# Main title block
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.8),
         "aic_ts_suite", font_size=52, color=C_WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.65), Inches(11), Inches(0.6),
         "Modular Time-Series Forecasting Toolkit", font_size=28, color=C_ACCENT3)
add_text(slide, Inches(1), Inches(3.35), Inches(11), Inches(0.4),
         "CTO Technical Briefing  |  Analytics Engineering Team", font_size=16, color=C_MUTED)

# Accent line
add_shape(slide, Inches(1), Inches(3.9), Inches(3), Pt(3),
          fill_color=C_ACCENT1, shape_type=MSO_SHAPE.RECTANGLE)

# Key stats row
stats = [
    ("9", "Forecasting\nModels"),
    ("4", "Algorithm\nParadigms"),
    ("5", "Evaluation\nKPIs"),
    ("35", "Python\nModules"),
    ("~3K", "Lines of\nCode"),
]
x_start = Inches(1)
for i, (num, label) in enumerate(stats):
    x = x_start + Inches(i * 2.2)
    add_shape(slide, x, Inches(4.5), Inches(1.8), Inches(1.4), fill_color=C_CARD_BG, border_color=C_BORDER)
    add_text(slide, x, Inches(4.6), Inches(1.8), Inches(0.6),
             num, font_size=36, color=C_ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.15), Inches(1.8), Inches(0.6),
             label, font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER)

# Date & version
add_text(slide, Inches(1), Inches(6.5), Inches(5), Inches(0.3),
         "March 2026  |  v0.1.0  |  Python 3.10+", font_size=12, color=C_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Executive Summary", "What this toolkit delivers")

# Left panel - description
add_multi_text(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(2.5), [
    ("End-to-end time-series forecasting framework that unifies:", 14, C_LIGHT),
    ("", 6),
    ("Data ingestion from CSV, Excel, TimescaleDB, and Open-Meteo weather API", 12, C_LIGHT),
    ("Automated cleaning with fluent DataCleaner pipeline", 12, C_LIGHT),
    ("Feature engineering: lags, rolling stats, Fourier harmonics", 12, C_LIGHT),
    ("9 forecasting models across 4 paradigms", 12, C_LIGHT),
    ("5 standardised KPIs with model comparison leaderboard", 12, C_LIGHT),
    ("UUID-based end-to-end run traceability", 12, C_LIGHT),
    ("YAML-driven experiment orchestration with checkpointing", 12, C_LIGHT),
])

# Right panel - key differentiators cards
cards = [
    ("Model-Agnostic Comparison", "Unified ForecastResult interface\nfor all 9 algorithms", C_ACCENT1),
    ("Zero-to-Leaderboard", "Complete pipeline in 20 lines\nof Python code", C_ACCENT4),
    ("Reproducible Experiments", "YAML config + checkpointing\n+ resume capability", C_ACCENT2),
    ("Full Traceability", "UUID v4 correlation IDs flow\nthrough every module", C_ACCENT5),
    ("Modular Design", "Use only what you need\nEach package is independent", C_ACCENT3),
]

for i, (title, body, color) in enumerate(cards):
    y = Inches(1.4) + Inches(i * 1.05)
    add_shape(slide, Inches(6.8), y, Inches(5.8), Inches(0.9), fill_color=C_CARD_BG, border_color=C_BORDER)
    # accent dot
    add_shape(slide, Inches(7.0), y + Inches(0.2), Inches(0.12), Inches(0.12), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_text(slide, Inches(7.3), y + Inches(0.08), Inches(5), Inches(0.3),
             title, font_size=13, color=color, bold=True)
    add_text(slide, Inches(7.3), y + Inches(0.38), Inches(5), Inches(0.45),
             body, font_size=10, color=C_LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: ARCHITECTURE OVERVIEW (Canvas diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Architecture Overview", "Layered module architecture")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=JRc6HIQ6JCYY6Hgb25AjN,SeNIBmzVCsEGjPgApUiVAg")

# Draw 4-layer architecture diagram
layers = [
    ("EVALUATION & OUTPUT", C_ACCENT6, [
        ("evaluation/", "5 KPIs, Leaderboard\nTraceability"),
        ("viz/", "5 plot types\nMatplotlib"),
        ("display/", "Styled HTML\nMetric cards"),
    ]),
    ("FORECASTING ENGINE", C_ACCENT4, [
        ("univariate/", "AutoARIMA\nAutoETS, HW"),
        ("ml_models/", "XGBoost\nProphet"),
        ("neural/", "NHITS, MLP\nTimeGPT"),
        ("multivariate/", "VAR +\nGranger"),
    ]),
    ("TRANSFORM & FEATURES", C_ACCENT2, [
        ("signals/", "Log, Sqrt\nBox-Cox (lambda)"),
        ("features/", "Lags, Rolling Stats\nFourier Harmonics"),
    ]),
    ("DATA INGESTION", C_ACCENT1, [
        ("config/", "Singleton\nUUID v4"),
        ("connectivity/", "CSV, Excel\nTimescaleDB"),
        ("cleaning/", "DataCleaner\nFluent API"),
        ("weather/", "Open-Meteo\nAPI"),
    ]),
]

diagram_left = Inches(0.6)
diagram_width = Inches(12.1)
layer_height = Inches(1.2)
layer_gap = Inches(0.15)
y_start = Inches(1.5)

for li, (layer_name, accent, modules) in enumerate(layers):
    y = y_start + li * (layer_height + layer_gap)

    # Layer background
    add_shape(slide, diagram_left, y, diagram_width, layer_height,
              fill_color=C_CARD_BG, border_color=accent, border_width=Pt(1.5))

    # Layer label
    add_text(slide, diagram_left + Inches(0.2), y + Inches(0.05), Inches(3), Inches(0.3),
             layer_name, font_size=11, color=accent, bold=True)

    # Module boxes
    mod_count = len(modules)
    mod_width = Inches(2.6)
    total_mod_width = mod_count * mod_width + (mod_count - 1) * Inches(0.15)
    mod_x_start = diagram_left + (diagram_width - total_mod_width) / 2

    for mi, (mod_name, mod_desc) in enumerate(modules):
        mx = mod_x_start + mi * (mod_width + Inches(0.15))
        my = y + Inches(0.35)
        add_shape(slide, mx, my, mod_width, Inches(0.72),
                  fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_BORDER)
        add_text(slide, mx + Inches(0.1), my + Inches(0.02), mod_width - Inches(0.2), Inches(0.25),
                 mod_name, font_size=11, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, mx + Inches(0.1), my + Inches(0.28), mod_width - Inches(0.2), Inches(0.4),
                 mod_desc, font_size=9, color=C_MUTED, alignment=PP_ALIGN.CENTER)

# Down arrows between layers
for i in range(3):
    y_arrow = y_start + (i + 1) * (layer_height + layer_gap) - layer_gap + Pt(1)
    add_shape(slide, diagram_left + diagram_width / 2 - Inches(0.08), y_arrow - Inches(0.08),
              Inches(0.16), Inches(0.16), fill_color=C_ACCENT3,
              shape_type=MSO_SHAPE.DOWN_ARROW)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: DATA FLOW PIPELINE (Canvas diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Data Pipeline Flow", "End-to-end data journey from ingestion to output")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=I18_GDCG_l5pUvUtMCJi0,E6xdB-w-tVsSC-15W27bBg")

# Pipeline stages as connected boxes flowing top to bottom, then right
stages = [
    ("1", "Data Sources", "CSV / Excel / TimescaleDB\nOpen-Meteo Weather API", C_ACCENT1),
    ("2", "DataCleaner", ".load() -> .sanitize()\n.set_datetime_index() -> .result()", C_ACCENT1),
    ("3", "Train / Test Split", "train = series[:-horizon]\ntest = series[-horizon:]", C_ACCENT3),
    ("4", "Transforms & Features", "Log/Sqrt/Box-Cox transforms\nLags + Rolling + Fourier matrix", C_ACCENT2),
    ("5", "Forecasting Engine", "9 models across 4 paradigms\nReturns ForecastResult", C_ACCENT4),
    ("6", "Evaluation", "5 KPIs computed\nModelComparison leaderboard", C_ACCENT6),
]

# Layout: 2 rows x 3 columns
box_w = Inches(3.6)
box_h = Inches(2.2)
gap_x = Inches(0.4)
gap_y = Inches(0.3)
start_x = Inches(0.6)
start_y = Inches(1.5)

positions = []
for i in range(6):
    row = i // 3
    col = i % 3
    x = start_x + col * (box_w + gap_x)
    y = start_y + row * (box_h + gap_y)
    positions.append((x, y))

for i, (num, title, desc, color) in enumerate(stages):
    x, y = positions[i]

    # Box
    add_shape(slide, x, y, box_w, box_h, fill_color=C_CARD_BG, border_color=color, border_width=Pt(2))

    # Number badge
    add_pill(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.35), Inches(0.35),
             num, fill_color=color, font_size=14)

    # Title
    add_text(slide, x + Inches(0.65), y + Inches(0.15), box_w - Inches(0.8), Inches(0.35),
             title, font_size=16, color=color, bold=True)

    # Description
    add_text(slide, x + Inches(0.2), y + Inches(0.65), box_w - Inches(0.4), Inches(1.4),
             desc, font_size=12, color=C_LIGHT)

    # Arrow to next (horizontal within row)
    if i % 3 < 2 and i < 5:
        ax = x + box_w + Inches(0.05)
        ay = y + box_h / 2
        add_shape(slide, ax, ay - Inches(0.07), Inches(0.3), Inches(0.14),
                  fill_color=C_ACCENT3, shape_type=MSO_SHAPE.RIGHT_ARROW)

# Down arrows between rows
for col in range(3):
    if col < 3:
        x = start_x + col * (box_w + gap_x) + box_w / 2 - Inches(0.07)
        # Only the last of top row needs down arrow
        if col == 2:
            y_a = start_y + box_h + Inches(0.03)
            add_shape(slide, x, y_a, Inches(0.14), Inches(0.24),
                      fill_color=C_ACCENT3, shape_type=MSO_SHAPE.DOWN_ARROW)

# Traceability note at bottom
add_shape(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
          fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_ACCENT5, border_width=Pt(1.5))
add_text(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5),
         "Traceability: UUID v4 correlation ID propagates through every stage: "
         "DataCleaner -> Forecaster -> ForecastResult -> ModelComparison -> AnalyticsEngineClient -> JSON",
         font_size=11, color=C_ACCENT5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: FORECASTING ALGORITHM COMPARISON
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Forecasting Algorithms", "9 models across 4 paradigms")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=tDY3C33kFr8SUdz-bOU9Y,4kzK6aF-nYfq8lFbp1GMlA")

# Paradigm columns
paradigms = [
    ("Classical", C_ACCENT1, [
        ("AutoARIMA", "Automatic SARIMA\norder selection\nvia AICc"),
        ("AutoETS", "18-combo grid\nError/Trend/Season\nselect by AIC"),
        ("Holt-Winters", "Exponential\nsmoothing, add/\nmul seasonality"),
    ]),
    ("Machine Learning", C_ACCENT4, [
        ("XGBoost", "Gradient-boosted\ntrees, 3 quantile\nmodels, recursive"),
        ("Prophet", "Trend + season +\nholidays, change-\npoint detection"),
    ]),
    ("Multivariate", C_ACCENT2, [
        ("VAR", "Vector Auto-\nRegression +\nGranger causality"),
    ]),
    ("Neural / Foundation", C_ACCENT3, [
        ("NHITS", "Hierarchical\ninterpolation\ndeep learning"),
        ("MLP", "Feed-forward\nneural network\nbaseline"),
        ("TimeGPT", "Foundation model\nZERO-SHOT\nNo training needed"),
    ]),
]

col_w = Inches(2.95)
col_gap = Inches(0.15)
col_start = Inches(0.6)
col_top = Inches(1.5)

for pi, (p_name, p_color, models) in enumerate(paradigms):
    x = col_start + pi * (col_w + col_gap)

    # Paradigm header
    add_shape(slide, x, col_top, col_w, Inches(0.4), fill_color=p_color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, x, col_top + Inches(0.05), col_w, Inches(0.3),
             p_name, font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Model cards
    for mi, (m_name, m_desc) in enumerate(models):
        my = col_top + Inches(0.55) + mi * Inches(1.55)
        add_shape(slide, x, my, col_w, Inches(1.4), fill_color=C_CARD_BG, border_color=C_BORDER)
        add_text(slide, x + Inches(0.1), my + Inches(0.08), col_w - Inches(0.2), Inches(0.3),
                 m_name, font_size=13, color=p_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), my + Inches(0.42), col_w - Inches(0.2), Inches(0.85),
                 m_desc, font_size=10, color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: ALGORITHM DEEP COMPARISON TABLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Algorithm Comparison Matrix", "Training, speed, prediction intervals, and trade-offs")

# Table data
headers = ["Model", "Paradigm", "Training", "Pred. Intervals", "Best For", "Key Limitation"]
rows = [
    ["AutoARIMA",    "Classical",   "seconds",     "Native (analytical)", "Seasonal + known period", "Linearity assumed"],
    ["AutoETS",      "Classical",   "seconds",     "Native (state space)","Smooth trends + season",  "No exogenous vars"],
    ["Holt-Winters", "Classical",   "seconds",     "Simulation-based",   "Stable seasonality",      "Fixed pattern"],
    ["VAR",          "Multivariate","seconds",     "Impulse-response",   "Cross-variable lags",     "All series stationary"],
    ["XGBoost",      "ML",          "sec-min",     "Quantile regression","Non-linear + features",   "Compounding errors"],
    ["Prophet",      "ML",          "minutes",     "Posterior sampling", "Trend breaks + holidays", "Slower training"],
    ["NHITS",        "Neural",      "min-hours",   "Not provided",       "Long-horizon complex",    "Needs lots of data"],
    ["MLP",          "Neural",      "min-hours",   "Not provided",       "Neural baseline",         "Needs lots of data"],
    ["TimeGPT",      "Foundation",  "ZERO-SHOT",   "API-provided",       "Cold-start / prototype",  "API dependency"],
]

row_colors = [C_ACCENT1, C_ACCENT1, C_ACCENT1, C_ACCENT2, C_ACCENT4, C_ACCENT4, C_ACCENT3, C_ACCENT3, C_ACCENT5]

# Draw table
table_left = Inches(0.4)
table_top = Inches(1.4)
col_widths = [Inches(1.5), Inches(1.4), Inches(1.2), Inches(2.2), Inches(2.8), Inches(3.2)]
row_h = Inches(0.5)
header_h = Inches(0.45)

# Header row
x = table_left
for j, h in enumerate(headers):
    add_shape(slide, x, table_top, col_widths[j], header_h,
              fill_color=RGBColor(0x25, 0x25, 0x40), shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, x + Inches(0.08), table_top + Inches(0.08), col_widths[j] - Inches(0.16), Inches(0.3),
             h, font_size=10, color=C_ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)
    x += col_widths[j]

# Data rows
for i, row in enumerate(rows):
    y = table_top + header_h + Inches(0.02) + i * (row_h + Inches(0.02))
    bg = C_CARD_BG if i % 2 == 0 else RGBColor(0x18, 0x18, 0x30)
    x = table_left
    for j, cell in enumerate(row):
        add_shape(slide, x, y, col_widths[j], row_h, fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
        c = row_colors[i] if j == 0 else C_LIGHT
        b = True if j == 0 else False
        fs = 10 if j > 0 else 11
        add_text(slide, x + Inches(0.08), y + Inches(0.1), col_widths[j] - Inches(0.16), Inches(0.3),
                 cell, font_size=fs, color=c, bold=b, alignment=PP_ALIGN.CENTER)
        x += col_widths[j]

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: MODEL DECISION FRAMEWORK (Canvas tree diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Model Selection Decision Tree", "When to use which algorithm")

# Decision tree as boxes with arrows
# Root question
add_shape(slide, Inches(4.5), Inches(1.5), Inches(4.3), Inches(0.9),
          fill_color=C_CARD_BG, border_color=C_ACCENT3, border_width=Pt(2))
add_text(slide, Inches(4.6), Inches(1.6), Inches(4.1), Inches(0.7),
         "Do you have training data\nand time to train?", font_size=14, color=C_ACCENT3,
         bold=True, alignment=PP_ALIGN.CENTER)

# No branch -> TimeGPT
add_shape(slide, Inches(0.8), Inches(3.0), Inches(2.8), Inches(0.9),
          fill_color=C_CARD_BG, border_color=C_ACCENT5, border_width=Pt(2))
add_text(slide, Inches(0.9), Inches(3.1), Inches(2.6), Inches(0.7),
         "TimeGPT\n(zero-shot, no training)", font_size=12, color=C_ACCENT5,
         bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.2), Inches(2.55), Inches(1), Inches(0.3),
         "NO", font_size=12, color=C_RED, bold=True, alignment=PP_ALIGN.CENTER)

# Yes branch -> Linear question
add_shape(slide, Inches(7.5), Inches(3.0), Inches(4.3), Inches(0.9),
          fill_color=C_CARD_BG, border_color=C_ACCENT3, border_width=Pt(2))
add_text(slide, Inches(7.6), Inches(3.1), Inches(4.1), Inches(0.7),
         "Is the relationship\nlinear & seasonal?", font_size=14, color=C_ACCENT3,
         bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.0), Inches(2.55), Inches(1), Inches(0.3),
         "YES", font_size=12, color=C_ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)

# Yes -> Classical
add_shape(slide, Inches(4.5), Inches(4.6), Inches(3.5), Inches(1.2),
          fill_color=C_CARD_BG, border_color=C_ACCENT1, border_width=Pt(2))
add_multi_text(slide, Inches(4.6), Inches(4.65), Inches(3.3), Inches(1.1), [
    ("Classical Models", 14, C_ACCENT1, True, PP_ALIGN.CENTER),
    ("AutoARIMA | AutoETS | Holt-Winters", 11, C_LIGHT, False, PP_ALIGN.CENTER),
    ("Fast, interpretable, proven", 10, C_MUTED, False, PP_ALIGN.CENTER),
])
add_text(slide, Inches(7.0), Inches(4.15), Inches(1), Inches(0.3),
         "YES", font_size=11, color=C_ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)

# No -> Multiple series question
add_shape(slide, Inches(9.0), Inches(4.6), Inches(3.8), Inches(0.9),
          fill_color=C_CARD_BG, border_color=C_ACCENT3, border_width=Pt(2))
add_text(slide, Inches(9.1), Inches(4.7), Inches(3.6), Inches(0.7),
         "Multiple related\ntime series?", font_size=14, color=C_ACCENT3,
         bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(11.5), Inches(4.15), Inches(1), Inches(0.3),
         "NO", font_size=11, color=C_RED, bold=True, alignment=PP_ALIGN.CENTER)

# Yes -> VAR
add_shape(slide, Inches(8.0), Inches(6.1), Inches(2.4), Inches(1.0),
          fill_color=C_CARD_BG, border_color=C_ACCENT2, border_width=Pt(2))
add_multi_text(slide, Inches(8.1), Inches(6.15), Inches(2.2), Inches(0.9), [
    ("VAR", 14, C_ACCENT2, True, PP_ALIGN.CENTER),
    ("Multivariate +", 10, C_LIGHT, False, PP_ALIGN.CENTER),
    ("Granger causality", 10, C_LIGHT, False, PP_ALIGN.CENTER),
])
add_text(slide, Inches(9.3), Inches(5.65), Inches(1), Inches(0.3),
         "YES", font_size=11, color=C_ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)

# No -> ML/Neural
add_shape(slide, Inches(10.8), Inches(6.1), Inches(2.2), Inches(1.0),
          fill_color=C_CARD_BG, border_color=C_ACCENT4, border_width=Pt(2))
add_multi_text(slide, Inches(10.9), Inches(6.15), Inches(2.0), Inches(0.9), [
    ("ML / Neural", 14, C_ACCENT4, True, PP_ALIGN.CENTER),
    ("XGBoost | Prophet", 10, C_LIGHT, False, PP_ALIGN.CENTER),
    ("NHITS | MLP", 10, C_LIGHT, False, PP_ALIGN.CENTER),
])
add_text(slide, Inches(11.8), Inches(5.65), Inches(1), Inches(0.3),
         "NO", font_size=11, color=C_RED, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: XGBoost DEEP DIVE (Canvas internals diagram)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "XGBoost Forecaster — Internal Architecture", "Feature engineering + triple-model quantile prediction")

# Step 1: Feature Matrix
add_shape(slide, Inches(0.4), Inches(1.5), Inches(3.8), Inches(5.2),
          fill_color=C_CARD_BG, border_color=C_ACCENT1, border_width=Pt(2))
add_pill(slide, Inches(0.6), Inches(1.6), Inches(1.0), Inches(0.35),
         "STEP 1", fill_color=C_ACCENT1, font_size=10)
add_text(slide, Inches(1.7), Inches(1.6), Inches(2.3), Inches(0.35),
         "Feature Matrix", font_size=16, color=C_ACCENT1, bold=True)

feat_items = [
    ("Lag Features", "lag_1 ... lag_12\nvia series.shift(k)", C_ACCENT3),
    ("Rolling Stats", "roll_mean_3/6/12\nroll_std_3/6/12", C_ACCENT2),
    ("Fourier Terms", "sin/cos pairs\nK harmonics at period", C_ACCENT5),
]
for fi, (fname, fdesc, fcolor) in enumerate(feat_items):
    fy = Inches(2.2) + fi * Inches(1.3)
    add_shape(slide, Inches(0.6), fy, Inches(3.4), Inches(1.1),
              fill_color=RGBColor(0x14, 0x14, 0x28), border_color=fcolor)
    add_text(slide, Inches(0.8), fy + Inches(0.08), Inches(3.0), Inches(0.3),
             fname, font_size=12, color=fcolor, bold=True)
    add_text(slide, Inches(0.8), fy + Inches(0.4), Inches(3.0), Inches(0.6),
             fdesc, font_size=10, color=C_LIGHT)

add_text(slide, Inches(0.6), Inches(5.8), Inches(3.4), Inches(0.5),
         "build_supervised_matrix()", font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER)

# Arrow
add_shape(slide, Inches(4.3), Inches(3.8), Inches(0.4), Inches(0.2),
          fill_color=C_ACCENT3, shape_type=MSO_SHAPE.RIGHT_ARROW)

# Step 2: Triple Model Training
add_shape(slide, Inches(4.9), Inches(1.5), Inches(3.8), Inches(5.2),
          fill_color=C_CARD_BG, border_color=C_ACCENT4, border_width=Pt(2))
add_pill(slide, Inches(5.1), Inches(1.6), Inches(1.0), Inches(0.35),
         "STEP 2", fill_color=C_ACCENT4, font_size=10)
add_text(slide, Inches(6.2), Inches(1.6), Inches(2.3), Inches(0.35),
         "3 XGBRegressors", font_size=16, color=C_ACCENT4, bold=True)

models_xgb = [
    ("Point Model", "reg:squarederror\nOptimal point forecast", C_ACCENT4),
    ("Lower Bound", "reg:quantileerror\nalpha/2 = 0.025", C_ACCENT1),
    ("Upper Bound", "reg:quantileerror\n1-alpha/2 = 0.975", C_ACCENT6),
]
for mi, (mname, mdesc, mcolor) in enumerate(models_xgb):
    my = Inches(2.2) + mi * Inches(1.3)
    add_shape(slide, Inches(5.1), my, Inches(3.4), Inches(1.1),
              fill_color=RGBColor(0x14, 0x14, 0x28), border_color=mcolor)
    add_text(slide, Inches(5.3), my + Inches(0.08), Inches(3.0), Inches(0.3),
             mname, font_size=12, color=mcolor, bold=True)
    add_text(slide, Inches(5.3), my + Inches(0.4), Inches(3.0), Inches(0.6),
             mdesc, font_size=10, color=C_LIGHT)

# Arrow
add_shape(slide, Inches(8.8), Inches(3.8), Inches(0.4), Inches(0.2),
          fill_color=C_ACCENT3, shape_type=MSO_SHAPE.RIGHT_ARROW)

# Step 3: Recursive Prediction
add_shape(slide, Inches(9.4), Inches(1.5), Inches(3.5), Inches(5.2),
          fill_color=C_CARD_BG, border_color=C_ACCENT3, border_width=Pt(2))
add_pill(slide, Inches(9.6), Inches(1.6), Inches(1.0), Inches(0.35),
         "STEP 3", fill_color=C_ACCENT3, font_size=10)
add_text(slide, Inches(10.7), Inches(1.6), Inches(2.0), Inches(0.35),
         "Recursive Predict", font_size=16, color=C_ACCENT3, bold=True)

add_multi_text(slide, Inches(9.6), Inches(2.3), Inches(3.1), Inches(4.0), [
    ("For each step t+1 to t+H:", 12, C_LIGHT, True),
    ("", 6),
    ("1. Compute features from\n   history (including prior\n   predictions)", 11, C_LIGHT),
    ("", 6),
    ("2. Predict with all 3 models:\n   point, lower, upper", 11, C_LIGHT),
    ("", 6),
    ("3. Append prediction to\n   history for next step", 11, C_LIGHT),
    ("", 6),
    ("4. Return ForecastResult\n   with forecast + PI bounds", 11, C_ACCENT4, True),
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: KPI METRICS DEEP DIVE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Evaluation KPIs", "5 metrics for comprehensive model assessment")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=IbDEGJdWuz1y2huDydQpH,mN27lFRAjLwPyxtUoNNZ_Q")

kpis = [
    ("MAE", "Mean Absolute Error",
     "(1/n) x SUM |y - y_hat|",
     "Average error in original units",
     "Simple, interpretable, robust to outliers",
     "All errors equally costly", C_ACCENT1),
    ("RMSE", "Root Mean Square Error",
     "SQRT( (1/n) x SUM (y - y_hat)^2 )",
     "Penalises large deviations more",
     "Differentiable, standard benchmark",
     "Large misses are costly", C_ACCENT4),
    ("MAPE", "Mean Abs Percentage Error",
     "(100/n) x SUM |y - y_hat| / |y|",
     "Scale-independent % error",
     "Cross-series comparison",
     "% accuracy across products", C_ACCENT5),
    ("sMAPE", "Symmetric MAPE",
     "(100/n) x SUM |y-y_hat| / ((|y|+|y_hat|)/2)",
     "Symmetric, handles near-zero",
     "Bounded 0-200%, fairer than MAPE",
     "Low-value data periods", C_ACCENT2),
    ("R-squared", "Coefficient of Determination",
     "1 - SS_res / SS_tot",
     "Proportion of variance explained",
     "Single quality score (1.0 = perfect)",
     "Overall model quality", C_ACCENT3),
]

card_w = Inches(2.35)
card_h = Inches(4.8)
card_gap = Inches(0.13)
card_start = Inches(0.45)
card_top = Inches(1.4)

for i, (abbr, full_name, formula, what, strength, use_when, color) in enumerate(kpis):
    x = card_start + i * (card_w + card_gap)

    # Card bg
    add_shape(slide, x, card_top, card_w, card_h, fill_color=C_CARD_BG, border_color=color, border_width=Pt(2))

    # Abbreviation header
    add_shape(slide, x, card_top, card_w, Inches(0.6), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, x, card_top + Inches(0.1), card_w, Inches(0.4),
             abbr, font_size=20, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Full name
    add_text(slide, x + Inches(0.1), card_top + Inches(0.7), card_w - Inches(0.2), Inches(0.35),
             full_name, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Formula
    add_shape(slide, x + Inches(0.1), card_top + Inches(1.15), card_w - Inches(0.2), Inches(0.65),
              fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_BORDER)
    add_text(slide, x + Inches(0.15), card_top + Inches(1.2), card_w - Inches(0.3), Inches(0.55),
             formula, font_size=8, color=C_ACCENT3, alignment=PP_ALIGN.CENTER)

    # What it measures
    add_text(slide, x + Inches(0.1), card_top + Inches(1.95), card_w - Inches(0.2), Inches(0.3),
             "MEASURES", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.1), card_top + Inches(2.2), card_w - Inches(0.2), Inches(0.55),
             what, font_size=9, color=C_LIGHT)

    # Strength
    add_text(slide, x + Inches(0.1), card_top + Inches(2.85), card_w - Inches(0.2), Inches(0.3),
             "STRENGTH", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.1), card_top + Inches(3.1), card_w - Inches(0.2), Inches(0.55),
             strength, font_size=9, color=C_LIGHT)

    # Use when
    add_text(slide, x + Inches(0.1), card_top + Inches(3.75), card_w - Inches(0.2), Inches(0.3),
             "USE WHEN", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.1), card_top + Inches(4.0), card_w - Inches(0.2), Inches(0.55),
             use_when, font_size=9, color=C_LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: METRIC SELECTION GUIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Metric Selection Guide", "Choosing the right KPI for your business question")

# Business question -> metric mapping
questions = [
    ("How far off are we in real units?", "MAE", "RMSE", C_ACCENT1,
     "Directly interpretable in the same units as your data. If your energy forecast is off by 50 MW on average, stakeholders immediately understand the impact."),
    ("How costly are our worst misses?", "RMSE", "MAE", C_ACCENT4,
     "Squares errors before averaging, so a single large miss (e.g. 500 MW spike) dominates. Critical for grid balancing and inventory where big errors are expensive."),
    ("What's our % accuracy across products?", "MAPE", "sMAPE", C_ACCENT5,
     "Scale-independent percentage allows comparing a 10-unit product vs a 10,000-unit product. Warning: unreliable near zero."),
    ("% accuracy with near-zero values?", "sMAPE", "MAE", C_ACCENT2,
     "Uses the average of actual and forecast in the denominator, avoiding division-by-zero. Bounded 0-200%. Preferred for intermittent demand."),
    ("Overall model quality score?", "R-squared", "RMSE", C_ACCENT3,
     "Single number from 0 to 1 showing how much variance the model captures. R-squared=0.92 means the model explains 92% of the data's variability."),
]

row_h = Inches(1.0)
row_gap = Inches(0.1)
table_top = Inches(1.4)
q_left = Inches(0.5)
q_width = Inches(3.5)
m_width = Inches(1.2)
desc_width = Inches(6.5)

for i, (question, primary, secondary, color, desc) in enumerate(questions):
    y = table_top + i * (row_h + row_gap)

    # Row background
    add_shape(slide, q_left, y, Inches(12.3), row_h, fill_color=C_CARD_BG, border_color=C_BORDER)

    # Question
    add_text(slide, q_left + Inches(0.15), y + Inches(0.15), q_width, Inches(0.7),
             f'"{question}"', font_size=11, color=C_WHITE, bold=True)

    # Primary metric pill
    add_pill(slide, q_left + q_width + Inches(0.2), y + Inches(0.25),
             Inches(1.0), Inches(0.35), primary, fill_color=color, font_size=11)

    # Secondary
    add_text(slide, q_left + q_width + Inches(1.35), y + Inches(0.3), Inches(0.8), Inches(0.3),
             f"+ {secondary}", font_size=10, color=C_MUTED)

    # Description
    add_text(slide, q_left + q_width + Inches(2.3), y + Inches(0.1), desc_width, Inches(0.85),
             desc, font_size=9, color=C_LIGHT)

# Default note
add_shape(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
          fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_ACCENT3)
add_text(slide, Inches(0.7), Inches(6.75), Inches(11.8), Inches(0.4),
         "Default leaderboard sort: RMSE (penalises large errors, most commonly used in production forecasting)",
         font_size=11, color=C_ACCENT3, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: CURRENT STATUS & READINESS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Current Toolkit Status", "Module readiness and codebase metrics")

# Module status table
modules = [
    ("config",           "1 file",   "2.3 KB",   "Production", C_ACCENT4, "Singleton, env-var driven, UUID"),
    ("connectivity",     "4 files",  "21.6 KB",  "Production", C_ACCENT4, "CSV, Excel, TimescaleDB, Weather"),
    ("cleaning",         "2 files",  "4.5 KB",   "Production", C_ACCENT4, "Fluent API, interpolate/ffill"),
    ("signals",          "1 file",   "5.6 KB",   "Production", C_ACCENT4, "Log, Sqrt, Box-Cox"),
    ("features",         "3 files",  "11.6 KB",  "Production", C_ACCENT4, "Lags, rolling, Fourier"),
    ("forecasting",      "5 files",  "41.4 KB",  "Production", C_ACCENT4, "9 models, 4 paradigms"),
    ("evaluation",       "3 files",  "10.1 KB",  "Production", C_ACCENT4, "5 KPIs, leaderboard, trace"),
    ("viz",              "5 files",  "12.1 KB",  "Production", C_ACCENT4, "5 plot types, styled HTML"),
    ("experiment runner", "1 file",  "~15 KB",   "Production", C_ACCENT4, "YAML-driven, checkpointing"),
]

t_headers = ["Module", "Files", "Size", "Status", "Description"]
t_col_w = [Inches(1.8), Inches(1.0), Inches(1.0), Inches(1.3), Inches(4.5)]
t_left = Inches(0.5)
t_top = Inches(1.4)
t_row_h = Inches(0.42)

# Headers
x = t_left
for j, h in enumerate(t_headers):
    add_shape(slide, x, t_top, t_col_w[j], t_row_h,
              fill_color=RGBColor(0x25, 0x25, 0x40), shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, x + Inches(0.08), t_top + Inches(0.07), t_col_w[j] - Inches(0.16), Inches(0.28),
             h, font_size=10, color=C_ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)
    x += t_col_w[j]

for i, (mod, files, size, status, scolor, desc) in enumerate(modules):
    y = t_top + t_row_h + Inches(0.02) + i * (t_row_h + Inches(0.02))
    bg = C_CARD_BG if i % 2 == 0 else RGBColor(0x18, 0x18, 0x30)
    vals = [mod, files, size, status, desc]
    x = t_left
    for j, val in enumerate(vals):
        add_shape(slide, x, y, t_col_w[j], t_row_h, fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
        c = C_ACCENT1 if j == 0 else (scolor if j == 3 else C_LIGHT)
        b = j == 0
        add_text(slide, x + Inches(0.08), y + Inches(0.07), t_col_w[j] - Inches(0.16), Inches(0.28),
                 val, font_size=9, color=c, bold=b, alignment=PP_ALIGN.CENTER if j < 4 else PP_ALIGN.LEFT)
        x += t_col_w[j]

# Summary stats on the right side
stats_data = [
    ("35", "Python Modules"),
    ("~3K", "Lines of Code"),
    ("9", "Forecast Models"),
    ("5", "KPI Metrics"),
    ("4", "Data Sources"),
    ("6", "Design Patterns"),
]

for i, (val, label) in enumerate(stats_data):
    sx = Inches(10.2)
    sy = Inches(1.4) + i * Inches(0.85)
    add_shape(slide, sx, sy, Inches(2.7), Inches(0.75), fill_color=C_CARD_BG, border_color=C_BORDER)
    add_text(slide, sx + Inches(0.1), sy + Inches(0.1), Inches(0.8), Inches(0.5),
             val, font_size=22, color=C_ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx + Inches(0.95), sy + Inches(0.2), Inches(1.6), Inches(0.35),
             label, font_size=10, color=C_LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: RISK & GAP ASSESSMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Risk & Gap Assessment", "Identified risks and mitigation recommendations")

risks = [
    ("Core Functionality", "LOW", C_ACCENT4,
     "All 9 modules are production-ready and fully functional",
     "Continue maintenance and feature development"),
    ("Test Coverage", "HIGH", C_RED,
     "0% test coverage - no pytest suite, no unit or integration tests",
     "Add pytest with unit tests for metrics, transforms, features + integration tests for forecasting pipeline"),
    ("CI/CD Pipeline", "HIGH", C_RED,
     "No GitHub Actions, no automated linting, testing, or build",
     "Add GitHub Actions workflow: ruff lint + pytest + build on every PR"),
    ("Documentation", "LOW", C_ACCENT4,
     "Comprehensive README, ONBOARDING guide, and ARCHITECTURE docs",
     "Maintain as features are added"),
    ("Packaging", "LOW", C_ACCENT4,
     "setup.py with optional extras [ml], [neural], [all]",
     "Consider migrating to pyproject.toml (modern standard)"),
    ("Deployment", "MEDIUM", C_ACCENT5,
     "K8s YAML exists but Docker Compose is placeholder only",
     "Build production Docker image, add Helm chart"),
]

card_w_risk = Inches(6.0)
card_h_risk = Inches(1.5)
left_col = Inches(0.5)
right_col = Inches(6.7)

for i, (area, level, lcolor, desc, rec) in enumerate(risks):
    col = left_col if i < 3 else right_col
    row = i if i < 3 else i - 3
    y = Inches(1.5) + row * (card_h_risk + Inches(0.12))

    add_shape(slide, col, y, card_w_risk, card_h_risk, fill_color=C_CARD_BG, border_color=C_BORDER)

    # Area name + risk badge
    add_text(slide, col + Inches(0.15), y + Inches(0.08), Inches(3), Inches(0.3),
             area, font_size=13, color=C_WHITE, bold=True)
    add_pill(slide, col + card_w_risk - Inches(1.1), y + Inches(0.08), Inches(0.9), Inches(0.3),
             level, fill_color=lcolor, font_size=10)

    # Description
    add_text(slide, col + Inches(0.15), y + Inches(0.45), card_w_risk - Inches(0.3), Inches(0.4),
             desc, font_size=9, color=C_LIGHT)

    # Recommendation
    add_text(slide, col + Inches(0.15), y + Inches(0.88), card_w_risk - Inches(0.3), Inches(0.2),
             "Recommendation:", font_size=8, color=C_ACCENT3, bold=True)
    add_text(slide, col + Inches(0.15), y + Inches(1.05), card_w_risk - Inches(0.3), Inches(0.4),
             rec, font_size=9, color=C_ACCENT3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: DESIGN PATTERNS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Design Patterns", "6 patterns powering the toolkit architecture")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=NXa2eVvbfBptGZGBFDNb9,zXs85MnAvQe5D28fAinTtQ")

patterns = [
    ("Fluent Builder", "DataCleaner", ".load().sanitize().result()", "Readable, chainable\ndata pipelines", C_ACCENT1),
    ("Strategy", "BaseTransform", "Log / Sqrt / BoxCox", "Pluggable transforms\nwithout code changes", C_ACCENT2),
    ("Abstract Factory", "BaseForecaster", "9 concrete models", "Drop-in model swap\nUniform API", C_ACCENT4),
    ("Singleton", "CONFIG", "Global state + UUID", "Single source of truth\nConsistent config", C_ACCENT5),
    ("Visitor/Aggregator", "ModelComparison", ".add().leaderboard()", "Heterogeneous model\nevaluation", C_ACCENT6),
    ("Correlation ID", "All Modules", "UUID v4 propagation", "End-to-end audit\ntrail across logs", C_ACCENT3),
]

pat_w = Inches(3.9)
pat_h = Inches(2.4)
pat_gap = Inches(0.2)
pat_start_x = Inches(0.5)
pat_start_y = Inches(1.5)

for i, (name, where, impl, value, color) in enumerate(patterns):
    col = i % 3
    row = i // 3
    x = pat_start_x + col * (pat_w + pat_gap)
    y = pat_start_y + row * (pat_h + pat_gap)

    add_shape(slide, x, y, pat_w, pat_h, fill_color=C_CARD_BG, border_color=color, border_width=Pt(2))

    # Pattern name
    add_text(slide, x + Inches(0.15), y + Inches(0.1), pat_w - Inches(0.3), Inches(0.35),
             name, font_size=16, color=color, bold=True)

    # Where
    add_text(slide, x + Inches(0.15), y + Inches(0.5), pat_w - Inches(0.3), Inches(0.2),
             "WHERE", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), pat_w - Inches(0.3), Inches(0.25),
             where, font_size=11, color=C_LIGHT)

    # Implementation
    add_text(slide, x + Inches(0.15), y + Inches(1.0), pat_w - Inches(0.3), Inches(0.2),
             "HOW", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(1.2), pat_w - Inches(0.3), Inches(0.25),
             impl, font_size=11, color=C_LIGHT)

    # Business value
    add_text(slide, x + Inches(0.15), y + Inches(1.55), pat_w - Inches(0.3), Inches(0.2),
             "VALUE", font_size=8, color=C_MUTED, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(1.75), pat_w - Inches(0.3), Inches(0.5),
             value, font_size=10, color=color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: EXPERIMENT ORCHESTRATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Experiment Orchestration", "YAML-driven, checkpoint-resumable experiment runner")
add_excalidraw_badge(slide, "https://excalidraw.com/#json=15v3fu0fGzV0OJ_AB2p7N,RZ74uM9kbSdaVENfq2swLw")

# Left: Flow diagram
flow_steps = [
    ("1", "Load params.yaml", "Validate configuration", C_ACCENT1),
    ("2", "Create experiment dir", "experiments/{name}_{ts}_{id}/", C_ACCENT3),
    ("3", "Iterate enabled models", "ARIMA, ETS, HW, XGB, Prophet, VAR", C_ACCENT4),
    ("4", "Fit & evaluate each", "Save checkpoint + run JSON", C_ACCENT2),
    ("5", "Export leaderboard", "leaderboard.json + summary.json", C_ACCENT5),
]

for i, (num, title, desc, color) in enumerate(flow_steps):
    y = Inches(1.6) + i * Inches(1.05)
    add_shape(slide, Inches(0.5), y, Inches(5.5), Inches(0.9),
              fill_color=C_CARD_BG, border_color=color, border_width=Pt(1.5))
    add_pill(slide, Inches(0.7), y + Inches(0.25), Inches(0.4), Inches(0.35),
             num, fill_color=color, font_size=12)
    add_text(slide, Inches(1.25), y + Inches(0.1), Inches(4.5), Inches(0.35),
             title, font_size=13, color=color, bold=True)
    add_text(slide, Inches(1.25), y + Inches(0.48), Inches(4.5), Inches(0.35),
             desc, font_size=10, color=C_LIGHT)

    # Down arrow
    if i < 4:
        add_shape(slide, Inches(3.1), y + Inches(0.92), Inches(0.14), Inches(0.1),
                  fill_color=C_ACCENT3, shape_type=MSO_SHAPE.DOWN_ARROW)

# Right: Output structure
add_shape(slide, Inches(6.5), Inches(1.6), Inches(6.3), Inches(5.5),
          fill_color=C_CARD_BG, border_color=C_BORDER)
add_text(slide, Inches(6.7), Inches(1.7), Inches(5.9), Inches(0.35),
         "Output Structure", font_size=15, color=C_ACCENT3, bold=True)

folders = [
    ("params.yaml", "Frozen config snapshot", C_MUTED),
    ("checkpoints/", "Top-N model .pkl files by metric", C_ACCENT1),
    ("runs/", "Per-model run records (.json)", C_ACCENT4),
    ("test_runs/", "Test-set evaluation results", C_ACCENT2),
    ("logs/", "Structured per-model log files", C_MUTED),
    ("plots/", "Forecast visualisation PNGs", C_ACCENT6),
    ("results/", "leaderboard.json + summary.json", C_ACCENT5),
]

for i, (folder, desc, color) in enumerate(folders):
    fy = Inches(2.2) + i * Inches(0.6)
    add_text(slide, Inches(6.9), fy, Inches(2.0), Inches(0.3),
             folder, font_size=12, color=color, bold=True)
    add_text(slide, Inches(9.0), fy, Inches(3.5), Inches(0.3),
             desc, font_size=10, color=C_LIGHT)

# Key features
features = [
    "Checkpoint-based resumption (--resume flag)",
    "Top-N best model retention, pruned by metric",
    "Sprint snapshots at configurable intervals",
    "Structured logging with timestamps per model",
]
add_text(slide, Inches(6.7), Inches(5.5), Inches(5.9), Inches(0.25),
         "Key Features:", font_size=11, color=C_ACCENT3, bold=True)
for i, feat in enumerate(features):
    add_text(slide, Inches(6.9), Inches(5.8) + i * Inches(0.28), Inches(5.7), Inches(0.25),
             f"  {feat}", font_size=10, color=C_LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15: TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Technology Stack", "Dependencies and their roles")

stack = [
    ("Core", ["pandas >= 2.0", "numpy >= 1.24", "scipy >= 1.10", "matplotlib >= 3.7"], C_ACCENT1,
     "Data manipulation, numerical computation, visualization"),
    ("Statistical", ["statsmodels >= 0.14", "pmdarima >= 2.0"], C_ACCENT3,
     "ARIMA, ETS, Holt-Winters, VAR, Granger causality"),
    ("Machine Learning", ["xgboost >= 2.0", "prophet >= 1.1"], C_ACCENT4,
     "Gradient boosting, trend decomposition with changepoints"),
    ("Neural", ["neuralforecast >= 1.6", "(PyTorch backend)"], C_ACCENT2,
     "NHITS and MLP deep learning forecasters"),
    ("Foundation", ["nixtla >= 0.5", "NIXTLA_API_KEY env var"], C_ACCENT5,
     "TimeGPT zero-shot forecasting via API"),
    ("Database", ["sqlalchemy >= 2.0", "psycopg2-binary >= 2.9"], C_ACCENT6,
     "TimescaleDB time-series database connectivity"),
    ("Weather", ["Open-Meteo API", "(no API key required)"], C_ACCENT3,
     "Free historical and forecast weather data"),
    ("Deployment", ["Kubernetes", "Docker Compose"], C_MUTED,
     "GPU cluster job execution, local dev environment"),
]

col_w_s = Inches(6.0)
row_h_s = Inches(0.65)
left_x = Inches(0.5)

for i, (layer, packages, color, purpose) in enumerate(stack):
    y = Inches(1.4) + i * (row_h_s + Inches(0.06))

    # Layer badge
    add_shape(slide, left_x, y, Inches(1.8), row_h_s, fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, left_x, y + Inches(0.15), Inches(1.8), Inches(0.35),
             layer, font_size=11, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Packages
    pkg_text = "  |  ".join(packages)
    add_shape(slide, left_x + Inches(1.85), y, Inches(4.5), row_h_s,
              fill_color=C_CARD_BG, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, left_x + Inches(2.0), y + Inches(0.15), Inches(4.2), Inches(0.35),
             pkg_text, font_size=10, color=C_LIGHT)

    # Purpose
    add_shape(slide, left_x + Inches(6.4), y, Inches(5.9), row_h_s,
              fill_color=C_CARD_BG, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, left_x + Inches(6.55), y + Inches(0.15), Inches(5.6), Inches(0.35),
             purpose, font_size=10, color=C_MUTED)

# Install commands
add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
          fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_ACCENT3)
add_text(slide, Inches(0.7), Inches(6.45), Inches(11.8), Inches(0.6),
         'pip install -e .          # Core only\n'
         'pip install -e ".[ml]"    # + XGBoost, Prophet\n'
         'pip install -e ".[all]"   # Everything',
         font_size=10, color=C_ACCENT3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16: ROADMAP & NEXT STEPS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title_bar(slide, "Roadmap & Recommendations", "Prioritised next steps for production readiness")

phases = [
    ("IMMEDIATE", "Critical", C_RED, [
        ("Add Test Suite", "pytest with unit + integration tests for all modules"),
        ("Set Up CI/CD", "GitHub Actions: ruff lint + pytest + build on every PR"),
    ]),
    ("SHORT-TERM", "High Value", C_ACCENT5, [
        ("Migrate to pyproject.toml", "Modern Python packaging standard"),
        ("Dockerize Toolkit", "Production-ready container image"),
        ("Add Backtesting", "Walk-forward cross-validation for robust evaluation"),
    ]),
    ("MEDIUM-TERM", "Strategic", C_ACCENT1, [
        ("Model Registry", "MLflow integration for experiment tracking at scale"),
        ("API Service Layer", "FastAPI wrapper for serving forecasts as REST endpoints"),
        ("Streaming Support", "Real-time forecast updates for live data feeds"),
    ]),
]

phase_w = Inches(3.9)
phase_h = Inches(5.0)
phase_gap = Inches(0.2)
phase_start = Inches(0.5)

for pi, (phase_name, priority, color, items) in enumerate(phases):
    x = phase_start + pi * (phase_w + phase_gap)
    y = Inches(1.5)

    # Phase container
    add_shape(slide, x, y, phase_w, phase_h, fill_color=C_CARD_BG, border_color=color, border_width=Pt(2))

    # Header
    add_shape(slide, x, y, phase_w, Inches(0.7), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, x, y + Inches(0.05), phase_w, Inches(0.3),
             phase_name, font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.35), phase_w, Inches(0.25),
             priority, font_size=10, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    # Items
    for ii, (item_title, item_desc) in enumerate(items):
        iy = y + Inches(0.9) + ii * Inches(1.3)
        add_shape(slide, x + Inches(0.15), iy, phase_w - Inches(0.3), Inches(1.1),
                  fill_color=RGBColor(0x14, 0x14, 0x28), border_color=C_BORDER)

        # Number
        add_pill(slide, x + Inches(0.25), iy + Inches(0.1), Inches(0.3), Inches(0.3),
                 str(ii + 1 + sum(len(p[3]) for p in phases[:pi])),
                 fill_color=color, font_size=10)

        add_text(slide, x + Inches(0.65), iy + Inches(0.08), phase_w - Inches(0.9), Inches(0.3),
                 item_title, font_size=12, color=color, bold=True)
        add_text(slide, x + Inches(0.3), iy + Inches(0.45), phase_w - Inches(0.6), Inches(0.55),
                 item_desc, font_size=10, color=C_LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17: CLOSING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=C_ACCENT1, shape_type=MSO_SHAPE.RECTANGLE)

add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
         "aic_ts_suite", font_size=48, color=C_WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.5),
         "Modular Time-Series Forecasting Toolkit", font_size=24, color=C_ACCENT3)

add_shape(slide, Inches(1), Inches(3.5), Inches(3), Pt(3),
          fill_color=C_ACCENT1, shape_type=MSO_SHAPE.RECTANGLE)

# Key takeaways
takeaways = [
    "9 forecasting models across 4 paradigms — from classical ARIMA to zero-shot TimeGPT",
    "5 standardised KPIs with automated model comparison leaderboard",
    "End-to-end data pipeline: ingestion, cleaning, features, forecasting, evaluation",
    "YAML-driven experiment orchestration with checkpoint/resume capability",
    "UUID v4 correlation IDs for full run traceability",
    "All modules production-ready — immediate priorities: testing & CI/CD",
]

for i, t in enumerate(takeaways):
    add_shape(slide, Inches(1.2), Inches(3.9) + i * Inches(0.45), Inches(0.1), Inches(0.1),
              fill_color=C_ACCENT1, shape_type=MSO_SHAPE.OVAL)
    add_text(slide, Inches(1.5), Inches(3.83) + i * Inches(0.45), Inches(10), Inches(0.4),
             t, font_size=13, color=C_LIGHT)

# Excalidraw links
add_text(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.35),
         "Interactive Architecture Diagrams: ARCHITECTURE.md  |  Analytics Engineering Team  |  March 2026",
         font_size=11, color=C_MUTED)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "/Users/karimmorad/Projects/TSAtoolkit/docs/CTO_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
